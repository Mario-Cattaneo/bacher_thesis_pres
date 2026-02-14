#!/usr/bin/env python3
"""Inspect a .potx/.pptx presentation and print slide contents summary."""
import argparse
import os
import shutil
import subprocess
import tempfile
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def summarize_presentation(path):
    prs = Presentation(path)
    print(f"File: {path}")
    print(f"Slide count: {len(prs.slides)}")
    print()

    for i, slide in enumerate(prs.slides, start=1):
        print(f"--- Slide {i} ---")
        # Title (if present)
        title_text = None
        try:
            title_shape = slide.shapes.title
            if title_shape is not None and getattr(title_shape, 'has_text_frame', False):
                title_text = title_shape.text.strip().replace('\n', ' ')
        except Exception:
            title_text = None
        if title_text:
            print(f"Title: {title_text}")

        layout = getattr(slide, 'slide_layout', None)
        if layout is not None:
            name = getattr(layout, 'name', None)
            if name:
                print(f"Layout: {name}")

        # Shapes
        for si, shape in enumerate(slide.shapes, start=1):
            stype = shape.shape_type
            typename = stype.name if hasattr(stype, 'name') else str(stype)
            info = f"Shape {si}: {typename}"
            # mark title shape
            try:
                if shape == getattr(slide.shapes, 'title', None):
                    info += ' [TITLE]'
            except Exception:
                pass
            if getattr(shape, 'is_placeholder', False):
                try:
                    pidx = shape.placeholder_format.idx
                    info += f" (placeholder idx={pidx})"
                except Exception:
                    pass
            if getattr(shape, 'has_text_frame', False):
                txt = shape.text.strip().replace('\n', ' ')
                snippet = txt[:300] + ("..." if len(txt) > 300 else "")
                info += f" | text='{snippet}'"
            # Picture detection
            if stype == MSO_SHAPE_TYPE.PICTURE or hasattr(shape, 'image'):
                try:
                    img = shape.image
                    blob = img.blob
                    info += f" | picture: {getattr(img, 'filename', 'embedded')} ({len(blob)} bytes, ext={img.ext})"
                except Exception:
                    info += " | picture: (could not read image)"

            print(info)

        # Notes
        try:
            notes_slide = slide.notes_slide
            notes_texts = []
            for nsh in notes_slide.shapes:
                if getattr(nsh, 'has_text_frame', False):
                    text = nsh.text.strip().replace('\n', ' ')
                    if text:
                        notes_texts.append(text)
            if notes_texts:
                notes_join = ' | '.join(notes_texts)[:400]
                print(f"Notes: {notes_join}")
        except Exception:
            pass

        # Related images (parts)
        imgs = []
        for rel in getattr(slide.part, 'related_parts', {}).values():
            ct = getattr(rel, 'content_type', '')
            if ct.startswith('image/'):
                try:
                    size = len(rel.blob)
                except Exception:
                    size = 'unknown'
                imgs.append(f"{rel.partname} ({ct}, {size} bytes)")
        if imgs:
            print("Images in slide:")
            for it in imgs:
                print(" -", it)

        print()


def main():
    parser = argparse.ArgumentParser(description='Inspect a .potx/.pptx file')
    parser.add_argument('file', nargs='?', default=None,
                        help='path to .pptx/.potx/.odp (defaults to disco-template.pptx|potx)')
    args = parser.parse_args()

    # choose a default file if none provided
    file_arg = args.file
    if file_arg is None:
        for cand in ('disco-template.pptx', 'disco-template.potx', 'disco-template.odp'):
            if os.path.exists(cand):
                file_arg = cand
                break
        if file_arg is None:
            parser.error('no input file provided and no disco-template.* found')

    # If input is .potx or .odp, convert to .pptx using soffice if available
    converted_tmpdir = None
    converted_path = file_arg
    ext = os.path.splitext(file_arg)[1].lower()
    needs_convert = ext in ('.potx', '.odp')
    if needs_convert and not file_arg.lower().endswith('.pptx'):
        soffice = shutil.which('soffice')
        if soffice:
            converted_tmpdir = tempfile.mkdtemp(prefix='pptx_convert_')
            try:
                subprocess.run([soffice, '--headless', '--convert-to', 'pptx', file_arg, '--outdir', converted_tmpdir], check=True)
                base = os.path.splitext(os.path.basename(file_arg))[0]
                candidate = os.path.join(converted_tmpdir, base + '.pptx')
                if os.path.exists(candidate):
                    converted_path = candidate
                    print(f'Converted {file_arg} -> {converted_path}')
                else:
                    print('Conversion produced no .pptx; will attempt to open original file')
            except subprocess.CalledProcessError:
                print('LibreOffice conversion failed; attempting to open original file')
        else:
            print('LibreOffice (soffice) not found; cannot auto-convert .potx/.odp')

    try:
        summarize_presentation(converted_path)
    finally:
        if converted_tmpdir:
            try:
                shutil.rmtree(converted_tmpdir)
            except Exception:
                pass


if __name__ == '__main__':
    main()
